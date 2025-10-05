import os
import base64
from typing import Optional, Any, List, Dict
from openai import OpenAI


class ParserOperations:
    """File parser operations class"""
    
    def __init__(self, openai_api_key: Optional[str] = None, openai_base_url: Optional[str] = None):
        """Initialize the parser operations handler.

        Args:
            openai_api_key: OpenAI API key. If not provided, reads from OPENAI_API_KEY environment variable.
            openai_base_url: OpenAI API base URL. If not provided, uses OPENAI_BASE_URL environment variable
                           or defaults to 'https://api.openai.com/v1'.
        """
        self.api_key = openai_api_key or os.getenv('OPENAI_API_KEY')
        self.base_url = openai_base_url or os.getenv('OPENAI_BASE_URL', 'https://api.openai.com/v1')
        self.model = os.getenv('OPENAI_MODEL', 'gpt-4o')
        self.max_tokens = int(os.getenv('MAX_TOKENS', '4096'))
        
        # initialize the OpenAI client
        if self.api_key:
            self.client = OpenAI(
                api_key=self.api_key,
                base_url=self.base_url
            )
        else:
            self.client = None
    
    def _base64_img(self, file_path: str) -> str:
        """Convert image file to base64 encoding.

        Args:
            file_path: Path to the image file to convert.

        Returns:
            str: Base64 encoded string of the image file.
        """
        with open(file_path, 'rb') as image_file:
            encoded_image = base64.b64encode(image_file.read()).decode('utf-8')
        return encoded_image
    
    def _base64_video(self, file_path: str, frame_interval: int = 10) -> List[str]:
        """Extract video frames and convert them to base64 encoding.

        Args:
            file_path: Path to the video file to process.
            frame_interval: Interval for frame extraction (every Nth frame).

        Returns:
            List[str]: List of base64 encoded frame images.
        """
        import cv2
        
        video = cv2.VideoCapture(file_path)
        base64_frames = []
        frame_count = 0
        while video.isOpened():
            success, frame = video.read()
            if not success:
                break
            if frame_count % frame_interval == 0:
                _, buffer = cv2.imencode('.jpg', frame)
                base64_frames.append(base64.b64encode(buffer).decode('utf-8'))
            frame_count += 1
        video.release()
        return base64_frames
    
    def _prepare_image_messages(self, task: str, base64_image: str) -> List[dict[str, Any]]:
        """prepare the image analysis message format"""
        return [
            {
                'role': 'user',
                'content': [
                    {'type': 'text', 'text': task},
                    {
                        'type': 'image_url',
                        'image_url': {'url': f'data:image/jpeg;base64,{base64_image}'},
                    },
                ],
            }
        ]
    
    def parse_pdf(self, file_path: str, save_path: str) -> Dict[str, Any]:
        """Parse a PDF file, extract text content and save to a file.
        
        Args:
            file_path: str: The path to the PDF file to parse.
            save_path: str: The path to save the parsed content.
            
        Returns:
            Dict[str, Any]: Dictionary containing parsing status and information.
        """
        if not os.path.exists(file_path):
            return {
                'success': False,
                'message': f'File not found: {file_path}',
            }
        
        if not save_path:
            return {
                'success': False,
                'message': 'Save path is required',
            }
        
        if os.path.exists(save_path):
            return {
                'success': False,
                'message': f'File already exists: {save_path}',
            }
        
        # check if the save_path is a directory
        if os.path.isdir(save_path):
            return {
                'success': False,
                'message': f'Save path is a directory, not a file: {save_path}',
            }
        
        # check if the parent directory of the save_path exists
        parent_dir = os.path.dirname(save_path)
        if parent_dir and not os.path.exists(parent_dir):
            return {
                'success': False,
                'message': f'Parent directory does not exist: {parent_dir}',
            }

        try:
            import PyPDF2
            
            result = f'[Reading PDF file from {file_path}]\n'
            content = PyPDF2.PdfReader(file_path)
            text = ''
            for page_idx in range(len(content.pages)):
                text += (
                    f'@@ Page {page_idx + 1} @@\n'
                    + content.pages[page_idx].extract_text()
                    + '\n\n'
                )
            result += text.strip()

            with open(save_path, 'w') as f:
                f.write(result)

            return {
                'success': True,
                'message': f'PDF file parsed and saved to `{save_path}`',
                'save_path': save_path,
            }
                
        except ImportError:
            return {
                'success': False,
                'message': 'Error: PyPDF2 library not installed. Please install it using: pip install PyPDF2',
            }
        except Exception as e:
            return {
                'success': False,
                'message': f'Error parsing PDF: {str(e)}',
            }
    
    def parse_docx(self, file_path: str, save_path: str) -> Dict[str, Any]:
        """Parse a DOCX file and save the parsed content to a file.
        
        Args:
            file_path: str: The path to the DOCX file to parse.
            save_path: str: The path to save the parsed content.
            
        Returns:
            Dict[str, Any]: Dictionary containing parsing status and information.
        """
        if not os.path.exists(file_path):
            return {
                'success': False,
                'message': f'File not found: {file_path}',
            }
        
        if not save_path:
            return {
                'success': False,
                'message': 'Save path is required',
            }
        
        if os.path.exists(save_path):
            return {
                'success': False,
                'message': f'File already exists: {save_path}',
            }
        
        # check if the save_path is a directory
        if os.path.isdir(save_path):
            return {
                'success': False,
                'message': f'Save path is a directory, not a file: {save_path}',
            }
        
        # check if the parent directory of the save_path exists
        parent_dir = os.path.dirname(save_path)
        if parent_dir and not os.path.exists(parent_dir):
            return {
                'success': False,
                'message': f'Parent directory does not exist: {parent_dir}',
            }
        
        try:
            import docx
            
            result = f'[Reading DOCX file from {file_path}]\n'
            content = docx.Document(file_path)
            text = ''
            for i, para in enumerate(content.paragraphs):
                text += f'@@ Page {i + 1} @@\n' + para.text + '\n\n'
            result += text

            with open(save_path, 'w') as f:
                f.write(result)

            return {
                'success': True,
                'message': f'DOCX file parsed and saved to {save_path}',
                'save_path': save_path,
            }
        except ImportError:
            return {
                'success': False,
                'message': 'Error: python-docx library not installed. Please install it using: pip install python-docx',
            }
        except Exception as e:
            return {
                'success': False,
                'message': f'Error parsing DOCX: {str(e)}',
            }
        
    def parse_latex(self, file_path: str, save_path: str) -> Dict[str, Any]:
        """Parse a LaTeX file and save the parsed content to a file.
        
        Args:
            file_path: str: The path to the LaTeX file to parse.
            save_path: str: The path to save the parsed content.
            
        Returns:
            Dict[str, Any]: Dictionary containing parsing status and information.
        """
        if not os.path.exists(file_path):
            return {
                'success': False,
                'message': f'File not found: {file_path}',
            }
        
        if not save_path:
            return {
                'success': False,
                'message': 'Save path is required',
            }
        
        
        if os.path.exists(save_path):
            return {
                'success': False,
                'message': f'File already exists: {save_path}',
            }
        
        # check if the save_path is a directory
        if os.path.isdir(save_path):
            return {
                'success': False,
                'message': f'Save path is a directory, not a file: {save_path}',
            }
        
        # check if the parent directory of the save_path exists
        parent_dir = os.path.dirname(save_path)
        if parent_dir and not os.path.exists(parent_dir):
            return {
                'success': False,
                'message': f'Parent directory does not exist: {parent_dir}',
            }
        
        try:
            from pylatexenc.latex2text import LatexNodes2Text
            
            result = f'[Reading LaTex file from {file_path}]\n'
            with open(file_path) as f:
                data = f.read()
            text = LatexNodes2Text().latex_to_text(data)
            result += text.strip()
            
            with open(save_path, 'w') as f:
                f.write(result)
            
            return {
                'success': True,
                'message': f'LaTeX file parsed and saved to {save_path}',
                'save_path': save_path,
            }
        except ImportError:
            return {
                'success': False,
                'message': 'Error: pylatexenc library not installed. Please install it using: pip install pylatexenc',
            }
        except Exception as e:
            return {
                'success': False,
                'message': f'Error parsing LaTeX: {str(e)}',
            }
        
    def parse_audio(self, file_path: str, save_path: str, model: str = 'whisper-1') -> Dict[str, Any]:
        """Parse an audio file, transcribe its content and save the parsed content to a file.
        
        Args:
            file_path: str: The path to the audio file to parse.
            save_path: str: The path to save the parsed content.
            model: str: The model to use for audio transcription. Default is 'whisper-1'.
            
        Returns:
            Dict[str, Any]: Dictionary containing parsing status and information.
        """
        if not os.path.exists(file_path):
            return {
                'success': False,
                'message': f'File not found: {file_path}',
            }
        
        if not save_path:
            return {
                'success': False,
                'message': 'Save path is required',
            }
        
        
        if os.path.exists(save_path):
            return {
                'success': False,
                'message': f'File already exists: {save_path}',
            }
        
        # check if the save_path is a directory
        if os.path.isdir(save_path):
            return {
                'success': False,
                'message': f'Save path is a directory, not a file: {save_path}',
            }
        
        # check if the parent directory of the save_path exists
        parent_dir = os.path.dirname(save_path)
        if parent_dir and not os.path.exists(parent_dir):
            return {
                'success': False,
                'message': f'Parent directory does not exist: {parent_dir}',
            }
        
        if not self.client:
            return {
                'success': False,
                'message': 'Error: OpenAI API key not provided. Please set OPENAI_API_KEY environment variable or provide it in constructor.',
            }
        
        result = f'[Transcribing audio file from {file_path}]\n'
        try:
            # call the OpenAI Whisper API to transcribe the audio file
            with open(file_path, 'rb') as audio_file:
                transcript = self.client.audio.translations.create(
                    model=model, 
                    file=audio_file
                )
            result += transcript.text

            with open(save_path, 'w') as f:
                f.write(result)

            return {
                'success': True,
                'message': f'Audio file transcribed and saved to {save_path}',
                'save_path': save_path,
            }
        except Exception as e:
            return {
                'success': False,
                'message': f'Error transcribing audio file: {str(e)}',
            }
    
    def parse_image(self, file_path: str, save_path: str, task: str = 'Describe this image as detail as possible.') -> Dict[str, Any]:
        """Parse an image file, analyze its content and save the parsed content to a file.
        
        Args:
            file_path: str: The path to the image file to parse.
            save_path: str: The path to save the parsed content.
            task: str: The task description for image analysis. Default is 'Describe this image as detail as possible.'
            
        Returns:
            Dict[str, Any]: Dictionary containing parsing status and information.
        """
        if not os.path.exists(file_path):
            return {
                'success': False,
                'message': f'File not found: {file_path}',
            }
        
        if not save_path:
            return {
                'success': False,
                'message': 'Save path is required',
            }
        
        if os.path.exists(save_path):
            return {
                'success': False,
                'message': f'File already exists: {save_path}',
            }
        
        # check if the save_path is a directory
        if os.path.isdir(save_path):
            return {
                'success': False,
                'message': f'Save path is a directory, not a file: {save_path}',
            }
        
        # check if the parent directory of the save_path exists
        parent_dir = os.path.dirname(save_path)
        if parent_dir and not os.path.exists(parent_dir):
            return {
                'success': False,
                'message': f'Parent directory does not exist: {parent_dir}',
            }
        
        if not self.client:
            return {
                'success': False,
                'message': 'Error: OpenAI API key not provided. Please set OPENAI_API_KEY environment variable or provide it in constructor.',
            }
        
        result = f'[Reading image file from {file_path}]\n'
        try:
            base64_image = self._base64_img(file_path)
            response = self.client.chat.completions.create(
                model=self.model,
                messages=self._prepare_image_messages(task, base64_image),
                max_tokens=self.max_tokens,
            )
            content = response.choices[0].message.content
            result += content

            with open(save_path, 'w') as f:
                f.write(result)

            return {
                'success': True,
                'message': f'Image file parsed and saved to {save_path}',
                'save_path': save_path,
            }
        except Exception as error:
            return {
                'success': False,
                'message': f'Error parsing image file: {str(error)}',
            }
    
    def parse_video(self, file_path: str, save_path: str, task: str = 'Describe this image as detail as possible.', frame_interval: int = 30) -> Dict[str, Any]:
        """Parse a video file, analyze its content and save the parsed content to a file.
        
        Args:
            file_path: str: The path to the video file to parse.
            save_path: str: The path to save the parsed content.
            task: str: The task description for video analysis. Default is 'Describe this image as detail as possible.'
            frame_interval: int: The frame interval for video analysis. Default is 30.
            
        Returns:
            Dict[str, Any]: Dictionary containing parsing status and information.
        """
        if not os.path.exists(file_path):
            return {
                'success': False,
                'message': f'File not found: {file_path}',
            }
        
        if not save_path:
            return {
                'success': False,
                'message': 'Save path is required',
            }
        
        if os.path.exists(save_path):
            return {
                'success': False,
                'message': f'File already exists: {save_path}',
            }
        
        # check if the save_path is a directory
        if os.path.isdir(save_path):
            return {
                'success': False,
                'message': f'Save path is a directory, not a file: {save_path}',
            }
        
        # check if the parent directory of the save_path exists
        parent_dir = os.path.dirname(save_path)
        if parent_dir and not os.path.exists(parent_dir):
            return {
                'success': False,
                'message': f'Parent directory does not exist: {parent_dir}',
            }
        
        if not self.client:
            return {
                'success': False,
                'message': 'Error: OpenAI API key not provided. Please set OPENAI_API_KEY environment variable or provide it in constructor.',
            }
        
        result = f'[Processing video file from {file_path} with frame interval {frame_interval}]\n'

        task = task or 'This is one frame from a video, please summarize this frame.'
        base64_frames = self._base64_video(file_path)
        selected_frames = base64_frames[::frame_interval]

        if len(selected_frames) > 30:
            new_interval = len(base64_frames) // 30
            selected_frames = base64_frames[::new_interval]

        result += f'Totally {len(selected_frames)} would be analyze...\n\n'

        idx = 0
        for base64_frame in selected_frames:
            idx += 1
            result += f'Process the {file_path}, current No. {idx * frame_interval} frame...\n'
            try:
                response = self.client.chat.completions.create(
                    model=self.model,
                    messages=self._prepare_image_messages(task, base64_frame),
                    max_tokens=self.max_tokens,
                )

                content = response.choices[0].message.content
                current_frame_content = f"Frame {idx}'s content: {content}\n"
                result += current_frame_content

            except Exception as error:
                result += f'Error with the request: {error}\n'
        
        with open(save_path, 'w') as f:
            f.write(result)

        return {
            'success': True,
            'message': f'Video file parsed and saved to {save_path}',
            'save_path': save_path,
        }
    
    def parse_pptx(self, file_path: str, save_path: str) -> Dict[str, Any]:
        """Parse a PPTX file and extract text content.
        
        Args:
            file_path: str: The path to the PPTX file to parse.
            save_path: str: The path to save the parsed content.
            
        Returns:
            Dict[str, Any]: Dictionary containing parsing status and information.
        """
        if not os.path.exists(file_path):
            return {
                'success': False,
                'message': f'File not found: {file_path}',
            }
        
        if not save_path:
            return{
                'success': False,
                'message': 'Save path is required',
            }
        
        if os.path.exists(save_path):
            return {
                'success': False,
                'message': f'File already exists: {save_path}',
            }
        
        # check if the save_path is a directory
        if os.path.isdir(save_path):
            return {
                'success': False,
                'message': f'Save path is a directory, not a file: {save_path}',
            }
        
        # check if the parent directory of the save_path exists
        parent_dir = os.path.dirname(save_path)
        if parent_dir and not os.path.exists(parent_dir):
            return {
                'success': False,
                'message': f'Parent directory does not exist: {parent_dir}',
            }
        
        try:
            from pptx import Presentation
            
            result = f'[Reading PowerPoint file from {file_path}]\n'
            pres = Presentation(str(file_path))
            text = []
            for slide_idx, slide in enumerate(pres.slides):
                text.append(f'@@ Slide {slide_idx + 1} @@')
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text.append(shape.text)
            result += '\n'.join(text)

            with open(save_path, 'w') as f:
                f.write(result)

            return {
                'success': True,
                'message': f'PowerPoint file parsed and saved to {save_path}',
                'save_path': save_path,
            }
        except ImportError:
            return {
                'success': False,
                'message': 'Error: python-pptx library not installed. Please install it using: pip install python-pptx',
            }
        except Exception as e:
            return {
                'success': False,
                'message': f'Error reading PowerPoint file: {str(e)}',
            }